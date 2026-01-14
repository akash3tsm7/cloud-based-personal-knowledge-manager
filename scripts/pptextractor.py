#!/usr/bin/env python
# -*- coding: utf-8 -*-
import sys
import os

def extract_pptx_text(pptx_path):
    """Extract text from PPTX file"""
    try:
        from pptx import Presentation
        
        if not os.path.exists(pptx_path):
            sys.stdout.write("ERROR: File not found\n")
            sys.stdout.flush()
            return False

        prs = Presentation(pptx_path)
        slides_text = []

        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text.strip())

            slide_content = " ".join(text_runs).strip()
            if slide_content:
                slides_text.append(f"Slide {i+1}: {slide_content}")

        # Output the text
        if slides_text:
            output = "\n\n".join(slides_text)
        else:
            output = "[No text found in presentation]"
        
        sys.stdout.write(output)
        sys.stdout.flush()
        return True
        
    except ImportError:
        sys.stdout.write("ERROR: python-pptx not installed. Run: pip install python-pptx\n")
        sys.stdout.flush()
        return False
    except Exception as e:
        sys.stdout.write(f"ERROR: {str(e)}\n")
        sys.stdout.flush()
        return False


if __name__ == "__main__":
    try:
        if len(sys.argv) < 2:
            sys.stdout.write("ERROR: No input file\n")
            sys.stdout.flush()
            sys.exit(1)

        pptx_path = sys.argv[1]
        success = extract_pptx_text(pptx_path)
        sys.exit(0 if success else 1)
        
    except Exception as e:
        sys.stdout.write(f"ERROR: Script error: {str(e)}\n")
        sys.stdout.flush()
        sys.exit(1)